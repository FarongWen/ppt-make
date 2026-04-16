#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
脚本功能：
1. 读取data.xlsx文件，对每行数据进行拼接（逗号分隔）
2. 读取ppt_template.pptx模版，替换文本框内容
3. 每页展示4或5行Excel数据（由输入参数决定）
4. 自动复制页面，直到所有Excel数据都填充完成
5. 保存为新的PPT文件

使用方法：
    python process_ppt.py [options]
    
选项：
    --excel EXCEL_FILE    指定Excel文件路径，默认为'data.xlsx'
    --template TEMPLATE   指定PowerPoint模板文件路径，默认为'ppt_template.pptx'
    --output OUTPUT       指定输出文件路径，默认为'output.pptx'
    --rows {4,5}          每页显示的行数，必须是4或5，默认为4

示例：
    python process_ppt.py --rows 5 --output custom_output.pptx
"""

import argparse
import pandas as pd
from pptx import Presentation
from pptx.util import Pt

def concat_row(row):
    """拼接一行的各个单元格数据，用逗号分隔"""
    # 过滤掉空值
    filtered_values = [str(val).strip() for val in row if pd.notna(val) and str(val).strip()]
    return "，".join(filtered_values)

def create_ppt(excel_file, template_file, output_file, rows_per_slide):
    """根据Excel内容创建PPT"""
    # 读取Excel文件
    df = pd.read_excel(excel_file)
    total_rows = len(df)
    
    # 读取PPT模板
    prs = Presentation(template_file)
    
    # 计算需要的幻灯片数量
    slide_count = (total_rows + rows_per_slide - 1) // rows_per_slide
    
    # 获取模板幻灯片上的文本框
    template_slide = prs.slides[0]
    
    # 查找文本框（通常是第二个形状，第一个是标题）
    text_shape = None
    for shape in template_slide.shapes:
        if hasattr(shape, 'text') and '陈芝琳' in shape.text:
            text_shape = shape
            break
    
    if not text_shape:
        raise ValueError("未在模板中找到目标文本框")
    
    # 处理第一张幻灯片
    rows_text = []
    for i in range(min(rows_per_slide, total_rows)):
        rows_text.append(concat_row(df.iloc[i]))
    text_shape.text = "\n\n".join(rows_text)
    
    # 为文本框设置字体属性
    for paragraph in text_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = '宋体'
            run.font.size = Pt(24)
    
    # 添加更多的幻灯片
    for slide_idx in range(1, slide_count):
        # 创建一个空白的幻灯片，不包含任何文本框
        blank_layout = prs.slide_layouts[6]  # 通常索引6是空白布局
        slide = prs.slides.add_slide(blank_layout)
        
        # # 复制标题
        # title_shape = None
        # # 查找标题形状 - 通常是第一个文本框
        # for shape in slide.shapes:
        #     if hasattr(shape, 'text_frame') and hasattr(shape, 'name') and '标题' in shape.name:
        #         title_shape = shape
        #         break
        #     # 尝试通过placeholder的类型来判断是否为标题
        #     elif hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
        #         # 通常标题的placeholder type为1
        #         if shape.placeholder_format.type == 1:
        #             title_shape = shape
        #             break
                
        # # 如果找不到标题，尝试使用第一个形状
        # if not title_shape and slide.shapes:
        #     for shape in slide.shapes:
        #         if hasattr(shape, 'text_frame'):
        #             title_shape = shape
        #             break
                    
        # # 设置标题文本和格式
        # if title_shape:
        #     title_shape.text = "发展党员工作议题"
            
        #     # 设置标题字体
        #     for paragraph in title_shape.text_frame.paragraphs:
        #         for run in paragraph.runs:
        #             run.font.name = '宋体'
        #             run.font.size = Pt(36)
        #             run.font.bold = True
        
        # 添加文本框
        start_row = slide_idx * rows_per_slide
        end_row = min(start_row + rows_per_slide, total_rows)
        
        # 获取文本框形状
        left = text_shape.left
        top = text_shape.top
        width = text_shape.width
        height = text_shape.height
        
        # 创建新的文本框
        textbox = slide.shapes.add_textbox(left, top, width, height)
        
        # 填充内容
        rows_text = []
        for i in range(start_row, end_row):
            rows_text.append(concat_row(df.iloc[i]))
        
        textbox.text_frame.text = "\n\n".join(rows_text)
        
        # 设置字体
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = '宋体'
                run.font.size = Pt(24)
    
    # 保存PPT
    prs.save(output_file)
    print(f"成功生成PPT文件：{output_file}")

def main():
    parser = argparse.ArgumentParser(description='根据Excel数据生成PPT')
    parser.add_argument('--excel', default='data.xlsx', help='Excel文件路径')
    parser.add_argument('--template', default='ppt_template.pptx', help='PPT模板文件路径')
    parser.add_argument('--output', default='output.pptx', help='输出PPT文件路径')
    parser.add_argument('--rows', type=int, default=4, choices=[3,4, 5], help='每页显示的行数（4或5）')
    
    args = parser.parse_args()
    
    create_ppt(args.excel, args.template, args.output, args.rows)

if __name__ == '__main__':
    main()
